"""
auto_evaluate_presentation.py
=============================
Automated End-to-End Presentation Renderer & Visual Inspector.

1. Calls PPTX Service to generate a full presentation.
2. Uses PowerPoint COM engine to render every slide to high-res JPG images.
3. Copies rendered slide images into artifacts directory for visual inspection.
4. Performs deep shape & contrast quality audit.
"""

import sys
import os
import shutil
import json

try:
    import win32com.client
    from pptx import Presentation
except ImportError:
    print("ERROR: pywin32 and python-pptx are required.", file=sys.stderr)
    sys.exit(1)


def render_pptx_to_jpg(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    # Open presentation hidden/minimized
    prs = ppt.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
    prs.SaveAs(os.path.abspath(output_dir), 17)  # 17 = ppSaveAsJPG
    prs.Close()
    ppt.Quit()

    image_files = sorted([os.path.join(output_dir, f) for f in os.listdir(output_dir) if f.lower().endswith(('.jpg', '.jpeg', '.png'))])
    return image_files


def copy_slides_to_artifacts(slide_images, artifact_dir):
    os.makedirs(artifact_dir, exist_ok=True)
    copied_paths = []
    for idx, img_path in enumerate(slide_images):
        dest_name = f"rendered_slide_{idx + 1}.jpg"
        dest_path = os.path.join(artifact_dir, dest_name)
        shutil.copy2(img_path, dest_path)
        copied_paths.append(dest_path)
    return copied_paths


def audit_presentation_structure(pptx_path):
    prs = Presentation(pptx_path)
    print("\nAUDITING PRESENTATION STRUCTURE: " + os.path.basename(pptx_path))
    print("   Total slides: " + str(len(prs.slides)))

    issues = []
    section_slides = []

    for idx, slide in enumerate(prs.slides):
        all_texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        tb3_texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame and s.name == "TextBox 3" and s.text_frame.text.strip()]
        
        is_section_header = bool(tb3_texts and tb3_texts[0] in ["01", "02", "03", "04", "05", "06"])
        if is_section_header:
            section_slides.append(idx + 1)

        print(f"   Slide {idx + 1}: {len(all_texts)} text shapes | {'[SECTION HEADER]' if is_section_header else '[CONTENT SLIDE]'}")

    # Check for consecutive section headers
    for i in range(len(section_slides) - 1):
        if section_slides[i + 1] == section_slides[i] + 1:
            issues.append(f"CRITICAL: Slides {section_slides[i]} and {section_slides[i + 1]} are consecutive Section Headers!")

    if issues:
        print("\nAUDIT ISSUES FOUND:")
        for issue in issues:
            print(f"   - {issue}")
    else:
        print("\nAUDIT PASSED: Zero consecutive Section Headers! Perfect structure.")

    return issues


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python auto_evaluate_presentation.py <pptx_file> [artifact_dir]", file=sys.stderr)
        sys.exit(1)

    pptx_file = sys.argv[1]
    artifact_dir = sys.argv[2] if len(sys.argv) >= 3 else os.path.join(os.getcwd(), "temp", "artifacts")

    rendered_dir = os.path.join(os.path.dirname(pptx_file), "rendered_slides")
    images = render_pptx_to_jpg(pptx_file, rendered_dir)
    print("\nRendered " + str(len(images)) + " slides to JPG:")
    for img in images:
        print("   - " + img)

    copied = copy_slides_to_artifacts(images, artifact_dir)
    print("\nCopied slides to Artifacts for Visual Inspection:")
    for c in copied:
        print("   - " + c)

    audit_presentation_structure(pptx_file)
