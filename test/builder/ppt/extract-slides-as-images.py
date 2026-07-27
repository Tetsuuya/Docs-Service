"""
Extract each slide from PPTX files as PNG images for visual inspection
Requires: python-pptx, Pillow, and potentially LibreOffice for rendering
"""
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("❌ Missing python-pptx. Installing...")
    os.system("pip install python-pptx")
    from pptx import Presentation

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("❌ Missing Pillow. Installing...")
    os.system("pip install Pillow")
    from PIL import Image, ImageDraw, ImageFont

def extract_pptx_info(pptx_path, output_dir):
    """Extract slide structure info and save as text-based visualization"""
    prs = Presentation(pptx_path)
    pptx_name = Path(pptx_path).stem
    
    output_folder = Path(output_dir) / pptx_name
    output_folder.mkdir(parents=True, exist_ok=True)
    
    print(f"\n📄 Analyzing: {pptx_name}")
    print(f"   {len(prs.slides)} slides\n")
    
    all_slides_info = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"   Slide {slide_idx}:")
        slide_info = {
            'slide_num': slide_idx,
            'shapes': []
        }
        
        # Analyze each shape on the slide
        for shape in slide.shapes:
            shape_info = {
                'type': shape.shape_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            
            # Extract text
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    shape_info['text'] = text[:100]  # First 100 chars
                    print(f"      📝 Text at ({shape.left/914400:.1f}, {shape.top/914400:.1f}): {text[:50]}...")
            
            # Check for images
            if hasattr(shape, "image"):
                shape_info['has_image'] = True
                print(f"      🖼️  Image at ({shape.left/914400:.1f}, {shape.top/914400:.1f})")
            
            # Check fill color
            if hasattr(shape, "fill"):
                if shape.fill.type == 1:  # Solid fill
                    try:
                        color = shape.fill.fore_color.rgb
                        shape_info['fill_color'] = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
                        print(f"      🎨 Fill: {shape_info['fill_color']}")
                    except:
                        pass
            
            slide_info['shapes'].append(shape_info)
        
        all_slides_info.append(slide_info)
        print()
    
    # Save structure to JSON
    import json
    with open(output_folder / 'structure.json', 'w') as f:
        json.dump(all_slides_info, f, indent=2)
    
    # Create visual text representation
    create_text_visualization(all_slides_info, output_folder, pptx_name)
    
    return all_slides_info

def create_text_visualization(slides_info, output_folder, pptx_name):
    """Create ASCII art visualization of slides"""
    
    report = []
    report.append("=" * 80)
    report.append(f"VISUAL LAYOUT ANALYSIS: {pptx_name}")
    report.append("=" * 80)
    report.append("")
    
    for slide in slides_info:
        report.append(f"\n{'─' * 80}")
        report.append(f"SLIDE {slide['slide_num']}")
        report.append(f"{'─' * 80}")
        
        # Group shapes by vertical position
        shapes_by_position = {}
        for shape in slide['shapes']:
            y_pos = int(shape['top'] / 914400)  # Convert to inches
            if y_pos not in shapes_by_position:
                shapes_by_position[y_pos] = []
            shapes_by_position[y_pos].append(shape)
        
        # Print in order from top to bottom
        for y_pos in sorted(shapes_by_position.keys()):
            shapes = shapes_by_position[y_pos]
            for shape in shapes:
                x_pos = int(shape['left'] / 914400)
                
                # Build shape description
                desc = f"  [{y_pos}in down, {x_pos}in right]"
                
                if 'text' in shape:
                    desc += f" TEXT: {shape['text'][:60]}"
                
                if 'has_image' in shape:
                    desc += f" [IMAGE]"
                
                if 'fill_color' in shape:
                    desc += f" Fill:{shape['fill_color']}"
                
                report.append(desc)
        
        report.append("")
    
    # Save report
    with open(output_folder / 'visual_layout.txt', 'w', encoding='utf-8') as f:
        f.write('\n'.join(report))
    
    print(f"✅ Saved visual layout report to {output_folder / 'visual_layout.txt'}")

def main():
    script_dir = Path(__file__).parent
    pptx_dir = script_dir / 'output'
    output_dir = script_dir / 'visual_analysis'
    
    if not pptx_dir.exists():
        print(f"❌ PPTX directory not found: {pptx_dir}")
        return
    
    pptx_files = list(pptx_dir.glob('test_*.pptx'))
    
    if not pptx_files:
        print(f"❌ No test_*.pptx files found in {pptx_dir}")
        return
    
    print("=" * 80)
    print("🔍 PPTX VISUAL LAYOUT ANALYZER")
    print("=" * 80)
    print(f"\nFound {len(pptx_files)} presentations to analyze\n")
    
    for pptx_file in pptx_files:
        try:
            extract_pptx_info(pptx_file, output_dir)
        except Exception as e:
            print(f"❌ Error processing {pptx_file.name}: {e}")
    
    print("\n" + "=" * 80)
    print("✅ Analysis complete!")
    print(f"📁 Results saved to: {output_dir}")
    print("=" * 80)

if __name__ == '__main__':
    main()
